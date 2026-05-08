#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
yt_slide_engine.py — YouTube AI動画 スライド生成エンジン

JSONスライドデータからWalkersブランドのPPTXを生成する。
テンプレートPPTXから冒頭5枚・CTA8枚をXMLコピーし、本編はpython-pptxで動的生成。

使い方:
  python3 yt_slide_engine.py --json slides.json --output slides.pptx
  python3 yt_slide_engine.py --json slides.json  # output は JSON の meta.slug から自動決定
"""

import argparse
import copy
import json
import os
import re
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── 定数 ─────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATE_PATH = SCRIPT_DIR / "template-slides.pptx"
CTA_IMAGES_DIR = SCRIPT_DIR / "cta-images"

# スライドサイズ (EMU)
SLIDE_W = 12191695
SLIDE_H = 6858000

# Walkersブランドカラー
CHARCOAL    = RGBColor(0x2B, 0x32, 0x3B)
ORANGE      = RGBColor(0xE9, 0x82, 0x12)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_ACCENT = RGBColor(0xFD, 0xF0, 0xDB)
LIGHT_BG    = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x66, 0x66, 0x66)
LIGHT_CHARCOAL = RGBColor(0x3D, 0x47, 0x53)
RED_ACCENT  = RGBColor(0xE0, 0x3E, 0x3E)
GREEN_ACCENT = RGBColor(0x2E, 0x8B, 0x57)
BLUE_ACCENT = RGBColor(0x3A, 0x7C, 0xBD)
PURPLE_ACCENT = RGBColor(0x7C, 0x3A, 0xED)

CARD_COLORS = {
    "":         ORANGE,
    "alt1":     BLUE_ACCENT,
    "alt2":     GREEN_ACCENT,
    "alt3":     PURPLE_ACCENT,
    "negative": RED_ACCENT,
}

FONT_NAME = "Yu Gothic"

# XML名前空間
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ─── ヘッダー・フッターの共通EMU値（テンプレートから抽出） ──────
HEADER_H     = Emu(1097280)
FOOTER_TOP   = Emu(6400800)
FOOTER_H     = Emu(457200)
CONTENT_TOP  = Emu(1371600)   # ヘッダー下のコンテンツ開始位置
CONTENT_BOTTOM = Emu(6400800)  # フッター上端
MARGIN_X     = Emu(731520)     # 左右マージン
CONTENT_W    = Emu(10515600)   # コンテンツ幅 = SLIDE_W - 2*MARGIN_X相当


# ════════════════════════════════════════════════════════
#  共通ヘルパー関数
# ════════════════════════════════════════════════════════

def set_font(run, name=FONT_NAME, size=None, bold=False, color=None):
    """runのフォントを設定する"""
    run.font.name = name
    if size:
        run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=Pt(18), bold=False, color=DARK_TEXT,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                word_wrap=True, space_after=None):
    """テキストボックスを追加する"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if space_after is not None:
            p.space_after = space_after
        run = p.add_run()
        run.text = line
        set_font(run, size=font_size, bold=bold, color=color)

    # 垂直位置（bodyPr要素のanchor属性に設定）
    body_pr = txBox.text_frame._txBody.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
    )
    if body_pr is not None:
        body_pr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))

    return txBox


def add_rect(slide, left, top, width, height, fill_color=CHARCOAL):
    """塗りつぶし矩形を追加する"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 枠線なし
    return shape


def new_blank_slide(prs):
    """空のBlankスライドを追加する"""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    # プレースホルダー除去
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


# ════════════════════════════════════════════════════════
#  XMLレベルコピー（テンプレートスライド用）
# ════════════════════════════════════════════════════════

def copy_slide(src_prs, src_idx, dst_prs):
    """テンプレートスライドをXMLレベルでコピーする（rId再マッピング付き）"""
    src_slide = src_prs.slides[src_idx]
    dst_slide = new_blank_slide(dst_prs)

    # Step 1: 画像リレーションを先にコピーし、旧rId→新rIdのマップを作る
    rid_map = {}
    for rel in src_slide.part.rels.values():
        if "image" in rel.reltype:
            result = dst_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            # python-pptx 1.x: get_or_add は rId 文字列を返す
            new_rId = result if isinstance(result, str) else result.rId
            rid_map[rel.rId] = new_rId

    # Step 2: Shape XMLをコピー
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        dst_slide.shapes._spTree.append(el)

    # Step 3: blip embed rId を新rIdに書き換え
    for blip in dst_slide._element.iter("{%s}blip" % NS_A):
        old_rId = blip.get("{%s}embed" % NS_R)
        if old_rId and old_rId in rid_map:
            blip.set("{%s}embed" % NS_R, rid_map[old_rId])

    return dst_slide


def copy_slide_replace_text(src_prs, src_idx, dst_prs, replacements):
    """テンプレートスライドをコピーし、テキストを差し替える

    replacements: dict of {shape_name: new_text}

    挙動:
    - テンプレの最初のパラグラフのフォント設定（size/bold/color/name）を全行に継承
    - 改行で複数行を入力した場合も全行で同じフォント
    - ①〜⑨ で始まる行は左寄せ + 左インデント（動画の趣旨スライドの内容リスト用）
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches as _Inches

    dst_slide = copy_slide(src_prs, src_idx, dst_prs)

    for shape in dst_slide.shapes:
        if shape.name in replacements and shape.has_text_frame:
            new_text = replacements[shape.name]
            tf = shape.text_frame

            # テンプレ最初のパラグラフからフォント設定を取得（全行に適用）
            template_font = {'name': None, 'size': None, 'bold': None, 'color': None}
            template_alignment = None
            if tf.paragraphs:
                src_p = tf.paragraphs[0]
                template_alignment = src_p.alignment
                if src_p.runs:
                    src_run = src_p.runs[0]
                    template_font['name'] = src_run.font.name
                    template_font['size'] = src_run.font.size
                    template_font['bold'] = src_run.font.bold
                    if src_run.font.color and src_run.font.color.type is not None:
                        try:
                            template_font['color'] = src_run.font.color.rgb
                        except AttributeError:
                            pass

            new_lines = new_text.split("\n")
            for i, line in enumerate(new_lines):
                if i < len(tf.paragraphs):
                    p = tf.paragraphs[i]
                else:
                    p = tf.add_paragraph()

                # ①〜⑨ で始まる行は左揃え + 左インデント
                # 「3行ブロックを左揃えで揃え、ブロック自体は画面中央寄りに見せる」設計。
                # 1.5 inch は手動修正版（見本）に近い経験値。
                if line and line[0] in '①②③④⑤⑥⑦⑧⑨':
                    p.alignment = PP_ALIGN.LEFT
                    # python-pptx の Paragraph には left_indent プロパティが無いので
                    # XML 属性 marL を直接設定（EMU単位、1 inch = 914400 EMU）
                    pPr = p._p.get_or_add_pPr()
                    pPr.set('marL', str(int(_Inches(2.1))))
                    pPr.set('indent', '0')
                else:
                    if template_alignment is not None:
                        p.alignment = template_alignment

                # runsを全削除
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)

                run = p.add_run()
                run.text = line
                if template_font['name']:
                    run.font.name = template_font['name']
                if template_font['size']:
                    run.font.size = template_font['size']
                if template_font['bold'] is not None:
                    run.font.bold = template_font['bold']
                if template_font['color']:
                    run.font.color.rgb = template_font['color']
                if not template_font['size']:
                    set_font(run, size=Pt(24), color=WHITE)

            # 余分なパラグラフを削除
            while len(tf.paragraphs) > len(new_lines):
                last_p = tf.paragraphs[-1]._p
                last_p.getparent().remove(last_p)

    return dst_slide


# ════════════════════════════════════════════════════════
#  共通レイアウトコンポーネント
# ════════════════════════════════════════════════════════

def add_header_bar(slide, section_num, header_text, footer_text):
    """ヘッダーバー（セクション番号 + タイトル）とフッターを追加する"""
    header_h = 1097280
    add_rect(slide, 0, 0, SLIDE_W, header_h, CHARCOAL)

    # 垂直中心を共有（タイトボックス + 位置計算で明示的に揃える）
    center_y = header_h // 2  # 548640

    # セクション番号
    num_box_h = int(Pt(44)) + 91440  # フォント + パディング
    add_textbox(slide, 548640, center_y - num_box_h // 2, 914400, num_box_h,
                section_num, font_size=Pt(44), bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ヘッダータイトル
    title_box_h = int(Pt(28)) + 91440  # フォント + パディング
    add_textbox(slide, 1645920, center_y - title_box_h // 2, 9144000, title_box_h,
                header_text, font_size=Pt(28), bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # フッター
    add_footer(slide, footer_text)


def add_footer(slide, footer_text):
    """フッターバーを追加する"""
    add_rect(slide, 0, 6400800, SLIDE_W, 457200, CHARCOAL)
    add_textbox(slide, 457200, 6446520, 10972800, 365760,
                footer_text, font_size=Pt(12), color=WHITE,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


DESC_HEIGHT = 365760  # 1行テキスト用（Pt22相当）

def add_description(slide, text, top=1371600):
    """ヘッダー下の説明テキスト（1行）"""
    add_textbox(slide, 731520, top, 10515600, DESC_HEIGHT,
                text, font_size=Pt(22), color=DARK_TEXT,
                alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════
#  本編スライド生成関数
# ════════════════════════════════════════════════════════

def make_section_title(prs, data):
    """セクションタイトルスライド（全面CHARCOAL + 番号 + タイトル）"""
    slide = new_blank_slide(prs)

    # 全面背景
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)

    # 区切り線
    add_rect(slide, 4572000, 2743200, 3047695, 54864, ORANGE)

    # セクション番号
    add_textbox(slide, 0, 1828800, SLIDE_W, 914400,
                data.get("num", "01"),
                font_size=Pt(48), bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # タイトル
    add_textbox(slide, 0, 3017520, SLIDE_W, 914400,
                data["title"],
                font_size=Pt(36), bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    return slide


def make_cards_slide(prs, data, footer_text):
    """カード型スライド（2〜3枚横並び）"""
    slide = new_blank_slide(prs)
    cards = data["cards"]
    n_cards = len(cards)
    header = data.get("header", "")
    section_num = data.get("num", "01")
    description = data.get("description", "")

    add_header_bar(slide, section_num, header, footer_text)

    # 説明文（オプション）- ヘッダー→desc→カードの余白を均等化
    if description:
        add_description(slide, description, top=1371600)
        card_top = 1920240
    else:
        card_top = 1371600

    # カード配置計算（コンテンツベース + 最低高さ保証）
    SUMMARY_BAR_H = 457200
    label_h = 457200
    max_items = max(len(c.get("items", [])) for c in cards)
    item_line_h = 350000
    card_h = label_h + max_items * item_line_h + 182880
    summary = data.get("summary", "")
    available_h = 6400800 - 182880 - card_top
    if summary:
        card_space = 6400800 - card_top - SUMMARY_BAR_H - 182880
        card_h = max(card_h, int(card_space * 0.85))
        card_h = min(card_h, card_space)
    else:
        card_h = max(card_h, int(available_h * 0.6))
        card_h = min(card_h, available_h)

    if n_cards == 2:
        card_w = 5029200
        positions = [731520, 6217920]
        body_inset_x = 137160
    elif n_cards == 3:
        card_w = 3474720
        positions = [457200, 4297680, 8138160]
        body_inset_x = 137160
    else:
        # 1枚カードの場合（稀）
        card_w = 10515600
        positions = [731520]
        body_inset_x = 137160

    for i, card in enumerate(cards):
        x = positions[i] if i < len(positions) else positions[-1]
        color_key = card.get("color", "")
        accent = CARD_COLORS.get(color_key, ORANGE)

        # カード背景
        add_rect(slide, x, card_top, card_w, card_h, LIGHT_BG)

        # ラベルバー
        add_rect(slide, x, card_top, card_w, label_h, accent)

        # ラベルテキスト
        label = card.get("label", "")
        add_textbox(slide, x + body_inset_x, card_top + 45720,
                    card_w - body_inset_x * 2, label_h - 91440,
                    label, font_size=Pt(20), bold=True, color=WHITE,
                    alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        # カード本文（手動垂直中央配置）
        items = card.get("items", [])
        body_text = "\n".join(items)
        n_items = len(items)

        # テキスト量に応じてフォントサイズを調整
        total_chars = sum(len(item) for item in items)
        if n_cards == 3:
            font_size = Pt(17) if total_chars > 60 else Pt(18)
        else:
            font_size = Pt(20) if total_chars > 80 else Pt(22)

        # テキスト高さを推定してカード内で中央に配置
        body_area_h = card_h - label_h
        text_block_h = int(n_items * font_size * 1.5)
        margin = max((body_area_h - text_block_h) // 2, 91440)
        body_top = card_top + label_h + margin
        body_h = body_area_h - margin * 2

        add_textbox(slide, x + body_inset_x, body_top,
                    card_w - body_inset_x * 2, body_h,
                    body_text, font_size=font_size, color=DARK_TEXT,
                    alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # 説明文がある場合は説明テキストを追加
    if description:
        add_description(slide, description)

    # サマリーバー（カード下〜フッター間を均等分割）
    if summary:
        bar_gap = (6400800 - card_top - card_h - SUMMARY_BAR_H) // 2
        s_top = card_top + card_h + bar_gap
        s_h = SUMMARY_BAR_H
        add_rect(slide, 457200, s_top, 11247120, s_h, LIGHT_ACCENT)
        add_textbox(slide, 731520, s_top,
                    10515600, s_h,
                    summary, font_size=Pt(18), color=DARK_TEXT,
                    alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                    space_after=Pt(0))

    return slide


def make_table_slide(prs, data, footer_text):
    """テーブル型スライド"""
    slide = new_blank_slide(prs)
    section_num = data.get("num", "01")
    header = data.get("header", "")
    columns = data.get("columns", [])
    rows = data.get("rows", [])
    description = data.get("description", "")

    add_header_bar(slide, section_num, header, footer_text)

    if description:
        add_description(slide, description)

    # テーブル位置
    table_top = 2286000 if description else 1645920
    table_left = 731520
    table_w = 10515600
    table_h = 6400800 - table_top - 457200 - 91440  # フッター上にマージン

    n_rows = len(rows) + 1  # ヘッダー行含む
    n_cols = len(columns)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                        Emu(table_left), Emu(table_top),
                                        Emu(table_w), Emu(table_h))
    table = tbl_shape.table

    # 列幅の設定
    if n_cols == 3:
        # #, 項目, 判定 のパターン
        col_widths = [Emu(int(table_w * 0.08)),
                      Emu(int(table_w * 0.72)),
                      Emu(int(table_w * 0.20))]
    elif n_cols == 2:
        col_widths = [Emu(int(table_w * 0.35)),
                      Emu(int(table_w * 0.65))]
    else:
        equal_w = Emu(table_w // n_cols)
        col_widths = [equal_w] * n_cols

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # 行の高さ（均等分割）
    row_h = Emu(table_h // n_rows)
    for i in range(n_rows):
        table.rows[i].height = row_h

    # ヘッダー行
    for j, col_name in enumerate(columns):
        cell = table.cell(0, j)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = CHARCOAL
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                set_font(run, size=Pt(18), bold=True, color=WHITE)

    # データ行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE

            # 判定列の特殊表示
            display_val = val
            if columns and j == len(columns) - 1 and val in ("yes", "no"):
                display_val = "○" if val == "yes" else "×"

            cell.text = display_val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    if val == "yes":
                        set_font(run, size=Pt(16), color=GREEN_ACCENT)
                    elif val == "no":
                        set_font(run, size=Pt(16), color=RED_ACCENT)
                    else:
                        set_font(run, size=Pt(16), color=DARK_TEXT)

    return slide


def make_comparison_slide(prs, data, footer_text):
    """比較型スライド（左右対比）"""
    slide = new_blank_slide(prs)
    section_num = data.get("num", "01")
    header = data.get("header", "")
    left = data.get("left", {})
    right = data.get("right", {})
    description = data.get("description", "")

    add_header_bar(slide, section_num, header, footer_text)

    if description:
        add_description(slide, description)
        card_top = 1920240
    else:
        card_top = 1645920
    card_w = 5029200
    SUMMARY_BAR_H = 457200
    label_h = 457200
    summary = data.get("summary", "")
    # コンテンツベース + 最低高さ保証
    max_items = max(len(left.get("items", [])), len(right.get("items", [])))
    item_line_h = 350000
    card_h = label_h + max_items * item_line_h + 182880
    available_h = 6400800 - 182880 - card_top
    if summary:
        card_space = 6400800 - card_top - SUMMARY_BAR_H - 182880
        card_h = max(card_h, int(card_space * 0.85))
        card_h = min(card_h, card_space)
    else:
        card_h = max(card_h, int(available_h * 0.6))
        card_h = min(card_h, available_h)

    # アイテム数に応じてフォント・行間を調整
    if max_items > 7:
        body_font = Pt(16)
        body_spacing = Pt(4)
    else:
        body_font = Pt(20)
        body_spacing = Pt(6)

    # 手動垂直中央配置の計算
    body_area_h = card_h - label_h

    def _comp_body_pos(items):
        n = len(items)
        text_h = int(n * (body_font + (body_spacing or Pt(0))) * 1.3)
        m = max((body_area_h - text_h) // 2, 91440)
        return card_top + label_h + m, body_area_h - m * 2

    # 左カード
    left_x = 731520
    left_color = CARD_COLORS.get(left.get("color", ""), ORANGE)
    add_rect(slide, left_x, card_top, card_w, card_h, LIGHT_BG)
    add_rect(slide, left_x, card_top, card_w, label_h, left_color)
    add_textbox(slide, left_x + 137160, card_top + 45720,
                card_w - 274320, label_h - 91440,
                left.get("label", ""), font_size=Pt(20), bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    left_items = left.get("items", [])
    l_top, l_h = _comp_body_pos(left_items)
    add_textbox(slide, left_x + 137160, l_top,
                card_w - 274320, l_h,
                "\n".join(left_items), font_size=body_font, color=DARK_TEXT,
                space_after=body_spacing, anchor=MSO_ANCHOR.MIDDLE)

    # 右カード
    right_x = 6217920
    right_color = CARD_COLORS.get(right.get("color", ""), BLUE_ACCENT)
    add_rect(slide, right_x, card_top, card_w, card_h, LIGHT_BG)
    add_rect(slide, right_x, card_top, card_w, label_h, right_color)
    add_textbox(slide, right_x + 137160, card_top + 45720,
                card_w - 274320, label_h - 91440,
                right.get("label", ""), font_size=Pt(20), bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    right_items = right.get("items", [])
    r_top, r_h = _comp_body_pos(right_items)
    add_textbox(slide, right_x + 137160, r_top,
                card_w - 274320, r_h,
                "\n".join(right_items), font_size=body_font, color=DARK_TEXT,
                space_after=body_spacing, anchor=MSO_ANCHOR.MIDDLE)

    # サマリーバー（カード下〜フッター間を均等分割）
    if summary:
        bar_gap = (6400800 - card_top - card_h - SUMMARY_BAR_H) // 2
        s_top = card_top + card_h + bar_gap
        s_h = SUMMARY_BAR_H
        add_rect(slide, 457200, s_top, 11247120, s_h, LIGHT_ACCENT)
        add_textbox(slide, 731520, s_top,
                    10515600, s_h,
                    summary, font_size=Pt(18), color=DARK_TEXT,
                    alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                    space_after=Pt(0))

    return slide


def make_flow_slide(prs, data, footer_text):
    """フロー型スライド（ステップ横並び）"""
    slide = new_blank_slide(prs)
    section_num = data.get("num", "01")
    header = data.get("header", "")
    steps = data.get("steps", [])
    description = data.get("description", "")

    add_header_bar(slide, section_num, header, footer_text)

    if description:
        add_description(slide, description)
        card_top = 1920240
    else:
        card_top = 1645920
    n_steps = len(steps)
    SUMMARY_BAR_H = 457200
    label_h = 457200
    summary = data.get("summary", "")

    # コンテンツベース + 最低高さ保証
    max_items = max(len(s.get("items", [])) for s in steps) if steps else 1
    item_line_h = 350000
    card_h = label_h + max_items * item_line_h + 182880
    available_h = 6400800 - 182880 - card_top
    if summary:
        card_space = 6400800 - card_top - SUMMARY_BAR_H - 182880
        card_h = max(card_h, int(card_space * 0.85))
        card_h = min(card_h, card_space)
    else:
        card_h = max(card_h, int(available_h * 0.6))
        card_h = min(card_h, available_h)

    # ステップカードの配置（カード型と同じロジック）
    gap = 274320
    total_gap = gap * (n_steps - 1)
    available_w = SLIDE_W - 914400  # 左右マージン
    card_w = (available_w - total_gap) // n_steps
    start_x = 457200

    for i, step in enumerate(steps):
        x = start_x + i * (card_w + gap)
        step_num = step.get("num", str(i + 1))
        label = step.get("label", f"STEP {step_num}")
        items = step.get("items", [])
        color_key = step.get("color", "")
        accent = CARD_COLORS.get(color_key, ORANGE)

        # カード背景
        add_rect(slide, x, card_top, card_w, card_h, LIGHT_BG)

        # ラベルバー
        add_rect(slide, x, card_top, card_w, label_h, accent)

        # ラベルテキスト
        add_textbox(slide, x + 91440, card_top + 45720,
                    card_w - 182880, label_h - 91440,
                    label, font_size=Pt(18), bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # ステップ内容（手動垂直中央配置）
        body_text = "\n".join(items)
        font_size = Pt(16) if n_steps >= 4 else Pt(18)
        n_items = len(items)
        flow_body_area_h = card_h - label_h
        flow_text_h = int(n_items * font_size * 1.5)
        flow_margin = max((flow_body_area_h - flow_text_h) // 2, 91440)
        add_textbox(slide, x + 91440, card_top + label_h + flow_margin,
                    card_w - 182880, flow_body_area_h - flow_margin * 2,
                    body_text, font_size=font_size, color=DARK_TEXT,
                    anchor=MSO_ANCHOR.MIDDLE)

        # 矢印（最後のステップ以外）
        if i < n_steps - 1:
            arrow_x = x + card_w + gap // 2 - 137160
            arrow_y = card_top + card_h // 2 - 137160
            add_textbox(slide, arrow_x, arrow_y, 274320, 274320,
                        "→", font_size=Pt(28), bold=True, color=ORANGE,
                        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # サマリーバー（カード下〜フッター間を均等分割）
    if summary:
        bar_gap = (6400800 - card_top - card_h - SUMMARY_BAR_H) // 2
        s_top = card_top + card_h + bar_gap
        s_h = SUMMARY_BAR_H
        add_rect(slide, 457200, s_top, 11247120, s_h, LIGHT_ACCENT)
        add_textbox(slide, 731520, s_top,
                    10515600, s_h,
                    summary, font_size=Pt(18), color=DARK_TEXT,
                    alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                    space_after=Pt(0))

    return slide


# ════════════════════════════════════════════════════════
#  メイン生成ロジック
# ════════════════════════════════════════════════════════

def generate_pptx(json_path, output_path=None, template_path=None):
    """JSONからPPTXを生成するメインエントリポイント"""

    # JSON読み込み
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    meta = data.get("meta", {})
    slides_data = data.get("slides", [])
    slug = meta.get("slug", "output")
    footer_text = meta.get("footer", f"{meta.get('title', '')} | アプリ開発研究所")

    # 出力パス決定
    if output_path is None:
        output_dir = SCRIPT_DIR.parent / slug
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "slides.pptx"

    # テンプレート読み込み
    if template_path is None:
        template_path = TEMPLATE_PATH
    tpl_prs = Presentation(str(template_path))

    # 出力用Presentation（テンプレートからサイズだけ継承）
    prs = Presentation()
    prs.slide_width = tpl_prs.slide_width
    prs.slide_height = tpl_prs.slide_height

    # テンプレートのスライド数確認
    tpl_count = len(tpl_prs.slides)
    print(f"Template: {tpl_count} slides")
    print(f"JSON: {len(slides_data)} slides defined")

    # ─── スライド生成 ───
    for i, sd in enumerate(slides_data):
        stype = sd.get("type", "")
        print(f"  Slide {i+1}: type={stype}")

        if stype == "title":
            # テンプレートスライド1からコピー、タイトルのみ差し替え
            # （タグ・サブタイトルは廃止：冒頭離脱率対策・yt-script SKILL.md 参照）
            title_text = sd.get("title", "")
            replacements = {}
            if title_text:
                replacements["TextBox 3"] = title_text
            copy_slide_replace_text(tpl_prs, 0, prs, replacements)

        elif stype == "text":
            # テンプレートスライド2（hero スタイル: 32pt Bold 中央配置）からコピー
            # 旧仕様で idx 4（チャンネル登録 24pt）を使っていたが、冒頭離脱率対策で
            # 太字・大きめフォントの hero テンプレに統一（2026-04-29 改修）
            lines = sd.get("lines", "")
            new_slide = copy_slide_replace_text(tpl_prs, 1, prs, {"TextBox 2": lines})
            # オプション: top_in 指定があれば TextBox の縦位置を上書き（単位: inch）
            top_in = sd.get("top_in")
            if top_in is not None:
                for sh in new_slide.shapes:
                    if sh.name == "TextBox 2" and sh.has_text_frame:
                        sh.top = int(float(top_in) * 914400)

        elif stype == "intro":
            # テンプレートスライド4をそのままコピー
            copy_slide(tpl_prs, 3, prs)

        elif stype == "subscribe":
            # テンプレートスライド5をそのままコピー
            copy_slide(tpl_prs, 4, prs)

        elif stype == "section":
            make_section_title(prs, sd)

        elif stype == "cards":
            make_cards_slide(prs, sd, footer_text)

        elif stype == "table":
            make_table_slide(prs, sd, footer_text)

        elif stype == "comparison":
            make_comparison_slide(prs, sd, footer_text)

        elif stype == "flow":
            make_flow_slide(prs, sd, footer_text)

        elif stype == "cta-text":
            # CTA テキストスライド — テンプレートからコピー
            # テンプレートの CTA スライドインデックスを特定
            cta_idx = _get_cta_index(sd, slides_data, i, tpl_count)
            if cta_idx is not None:
                copy_slide(tpl_prs, cta_idx, prs)
            else:
                # フォールバック: 動的生成
                _make_cta_text_fallback(prs, sd)

        elif stype == "cta-image":
            cta_idx = _get_cta_index(sd, slides_data, i, tpl_count)
            if cta_idx is not None:
                copy_slide(tpl_prs, cta_idx, prs)
            else:
                _make_cta_image_fallback(prs, sd)

        else:
            print(f"  WARNING: Unknown slide type '{stype}', skipping")

    # 保存
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    return str(output_path)


def _get_cta_index(sd, slides_data, current_idx, tpl_count):
    """CTA スライドのテンプレートインデックスを決定する

    テンプレートの最後8枚がCTA。
    slides_dataの最後8枚のうち、何番目かを算出してテンプレートインデックスに変換。
    """
    # slides_dataの末尾からCTAスライドをカウント
    cta_slides = []
    for j in range(len(slides_data) - 1, -1, -1):
        if slides_data[j].get("type", "").startswith("cta-"):
            cta_slides.insert(0, j)
        else:
            break

    if current_idx not in cta_slides:
        return None

    cta_position = cta_slides.index(current_idx)  # 0-based
    # テンプレートの最後8枚
    tpl_cta_start = tpl_count - 8  # 37-8=29 → index 29 (slide 30)
    tpl_idx = tpl_cta_start + cta_position

    if 0 <= tpl_idx < tpl_count:
        return tpl_idx
    return None


def _make_cta_text_fallback(prs, sd):
    """CTAテキストスライドの動的生成フォールバック"""
    slide = new_blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)
    text = sd.get("text", "")
    add_textbox(slide, 1371600, 1828800, 9448495, 3657600,
                text, font_size=Pt(24), color=WHITE,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def _make_cta_image_fallback(prs, sd):
    """CTA画像スライドの動的生成フォールバック"""
    slide = new_blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)
    src = sd.get("src", "")
    if src:
        img_path = CTA_IMAGES_DIR / Path(src).name
        if img_path.exists():
            from PIL import Image
            img = Image.open(str(img_path))
            aspect = img.width / img.height
            target_w = Emu(10000000)
            target_h = Emu(int(target_w / aspect))
            center_x = Emu((SLIDE_W - target_w) // 2)
            center_y = Emu((SLIDE_H - target_h) // 2)
            slide.shapes.add_picture(str(img_path), center_x, center_y, target_w, target_h)
    return slide


# ════════════════════════════════════════════════════════
#  検証
# ════════════════════════════════════════════════════════

def verify_pptx(pptx_path):
    """生成後のPPTXを検証する"""
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    errors = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                try:
                    blob_size = len(shape.image.blob)
                    if blob_size < 100:
                        errors.append(f"Slide {i+1}: image blob too small ({blob_size} bytes)")
                    else:
                        print(f"Slide {i+1}: image OK ({blob_size} bytes)")
                except Exception as e:
                    errors.append(f"Slide {i+1}: BROKEN image - {e}")

    if errors:
        print(f"\n*** FAIL: {len(errors)} broken images ***")
        for e in errors:
            print(f"  {e}")
        return False
    else:
        print(f"\nAll {total} slides verified. No broken images.")
        return True


# ════════════════════════════════════════════════════════
#  CLI
# ════════════════════════════════════════════════════════

_NUM_RE_SECTION = re.compile(r"^[0-9]{2}$")
_NUM_RE_CARDS = re.compile(r"^[①②③④⑤⑥⑦⑧⑨]$|^$")


def validate_slides_json(slides_data):
    """slides.json の構造を検証する。CRITICAL ルール違反があれば例外を投げる。

    検証項目:
      - section.num: '01'-'99' の2桁ゼロ埋めのみ
      - cards/comparison/table/flow.num: '①'-'⑨' または '' のみ（漢字・独自ラベル禁止）
      - section.title: 改行で分割した行数が 2 以下
    """
    errors = []
    for i, sd in enumerate(slides_data.get("slides", []), start=1):
        stype = sd.get("type")
        num = sd.get("num")

        if stype == "section":
            if num is not None and not _NUM_RE_SECTION.match(str(num)):
                errors.append(
                    f"slide {i} (section): num={num!r} は不正。"
                    f"section.num は '01'-'99' の2桁ゼロ埋めのみ許容。"
                    f"漢字（'結','序'等）や独自ラベル禁止。"
                )
            title = sd.get("title", "")
            line_count = len(title.split("\n"))
            if line_count > 2:
                errors.append(
                    f"slide {i} (section): title が {line_count} 行。"
                    f"section.title は 2 行以下に収めること（エンジンの textbox 高さは 1 inch 固定で 2 行までしか綺麗に収まらない）。"
                    f"現在: {title!r}"
                )
        elif stype in ("cards", "comparison", "table", "flow"):
            if num is not None and not _NUM_RE_CARDS.match(str(num)):
                errors.append(
                    f"slide {i} ({stype}): num={num!r} は不正。"
                    f"{stype}.num は '①'-'⑨' または '' のみ許容。"
                    f"漢字（'結','序'等）や数字（'01'）禁止。"
                )

    if errors:
        msg = "\n".join(f"  ❌ {e}" for e in errors)
        raise ValueError(
            f"slides.json 検証エラー（{len(errors)} 件）:\n{msg}\n"
            f"参照: .claude/skills/yt-slides/SKILL.md の「num フィールドの値域」「section title 行数制約」"
        )


def main():
    parser = argparse.ArgumentParser(description="YouTube AI動画 スライド生成エンジン")
    parser.add_argument("--json", required=True, help="JSONスライドデータファイルパス")
    parser.add_argument("--output", help="出力PPTXファイルパス（省略時: slugから自動決定）")
    parser.add_argument("--template", help="テンプレートPPTXパス（省略時: デフォルト）")
    parser.add_argument("--verify", action="store_true", help="生成後に検証を実行")
    parser.add_argument("--open", action="store_true", help="生成後にPPTXを開く")

    args = parser.parse_args()

    # JSON 構造を事前検証（CRITICAL ルール）
    with open(args.json, encoding="utf-8") as f:
        slides_data = json.load(f)
    validate_slides_json(slides_data)

    template = Path(args.template) if args.template else None
    output = Path(args.output) if args.output else None

    result_path = generate_pptx(args.json, output, template)

    if args.verify:
        ok = verify_pptx(result_path)
        if not ok:
            sys.exit(1)

    if args.open:
        os.system(f'open "{result_path}"')


if __name__ == "__main__":
    main()
